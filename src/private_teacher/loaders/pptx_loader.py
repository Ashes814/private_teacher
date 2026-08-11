"""
PPTX 加载器
===========

一张幻灯片 = 一个 Document（page = 幻灯片序号）。

python-pptx 的形状（shape）模型要点：
  - shape.has_text_frame 为 True 才有 .text（文本框、标题）
  - 表格是 shape.has_table，文字藏在 table.rows[i].cells[j].text 里，
    **不会**出现在 shape.text 中 —— 这是最常见的"PPT 内容丢了"的原因
  - 演讲者备注在 slide.notes_slide 里，也是有价值的内容
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from langchain_core.documents import Document
from loguru import logger

from private_teacher.loaders.base import base_metadata
from private_teacher.utils.exceptions import DocumentLoadError


def _shape_text(shape: Any) -> str:
    """从单个 shape 里抠出所有文字（含表格、组合形状）。"""
    parts: list[str] = []

    # ---------- 普通文本框 ----------
    if getattr(shape, "has_text_frame", False):
        text = shape.text_frame.text.strip()
        if text:
            parts.append(text)

    # ---------- 表格 ----------
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            # 一行拼成 "单元格1 | 单元格2"，保留表格的横向语义
            cells = [cell.text.strip() for cell in row.cells]
            line = " | ".join(c for c in cells if c)
            if line:
                parts.append(line)

    # ---------- 组合形状（递归） ----------
    # GROUP = 6，组合形状里还嵌套着子形状，不递归就会漏内容
    if getattr(shape, "shape_type", None) == 6:
        for sub in shape.shapes:
            sub_text = _shape_text(sub)
            if sub_text:
                parts.append(sub_text)

    return "\n".join(parts)


def load_pptx(
    path: Path,
    include_notes: bool = True,
    **_kwargs: Any,
) -> list[Document]:
    """加载 PPTX。

    Args:
        path: pptx 文件
        include_notes: 是否把演讲者备注也算进内容。
            默认 True —— 老师的备注往往比幻灯片正文信息量更大

    Returns:
        每张有内容的幻灯片一个 Document

    Raises:
        DocumentLoadError: 文件损坏 / 不是合法 pptx（比如老的 .ppt 格式）
    """
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
    except Exception as exc:
        raise DocumentLoadError(
            f"PPTX 解析失败（注意：python-pptx 不支持旧版 .ppt）: {exc}",
            path=str(path),
        ) from exc

    slides = list(prs.slides)
    total = len(slides)
    docs: list[Document] = []

    for slide_no, slide in enumerate(slides, start=1):
        parts = [_shape_text(shape) for shape in slide.shapes]

        # ---------- 演讲者备注 ----------
        if include_notes and slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                # 加个前缀，让 LLM 知道这段是备注而非正文
                parts.append(f"[演讲者备注] {notes}")

        content = "\n".join(p for p in parts if p).strip()
        if not content:
            continue  # 纯图片页

        docs.append(
            Document(
                page_content=content,
                metadata=base_metadata(path, page=slide_no, total_pages=total),
            )
        )

    logger.debug(f"{path.name}: 提取 {len(docs)}/{total} 张幻灯片")
    return docs
